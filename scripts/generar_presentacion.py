# -*- coding: utf-8 -*-
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def crear_presentacion():
    prs = Presentation()
    # Configurar formato panorámico 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colores corporativos oficiales (AEAT Estilo)
    AZUL_OSCURO = RGBColor(0, 68, 129)       # #004481
    AZUL_CLARO = RGBColor(238, 244, 250)     # #eef4fa
    TEXTO_OSCURO = RGBColor(30, 41, 59)      # #1e293b
    TEXTO_MUTED = RGBColor(100, 116, 139)    # #64748b
    AMBAR = RGBColor(217, 119, 6)            # #d97706
    BLANCO = RGBColor(255, 255, 255)
    VERDE = RGBColor(21, 128, 61)

    # Rutas de imágenes autogeneradas
    # Intenta buscar la imagen localmente o usa una ruta por defecto
    img_cadena = r"C:\Users\BRAIS\.gemini\antigravity\brain\abd2e431-f669-4796-a0bd-da6198329b9d\verifactu_chain_diagram_1780579995102.png"
    img_flujo = r"C:\Users\BRAIS\.gemini\antigravity\brain\abd2e431-f669-4796-a0bd-da6198329b9d\facturacion_flow_steps_1780580010392.png"

    # Layouts de pptx (en blanco para control de diseño)
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # DIAPOSITIVA 1: PORTADA
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = AZUL_OSCURO
    bg_shape.line.fill.background()

    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "LA NUEVA FACTURACIÓN EN ESPAÑA"
    p.font.bold = True
    p.font.size = Pt(42)
    p.font.color.rgb = BLANCO
    p.space_after = Pt(15)

    p2 = tf.add_paragraph()
    p2.text = "Guía Visual Sencilla para Autónomos y PYMEs"
    p2.font.bold = True
    p2.font.size = Pt(26)
    p2.font.color.rgb = AMBAR
    p2.space_after = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "Entendiendo Veri*factu, la Ley Crea y Crece, y el porqué de la digitalización fiscal."
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(147, 197, 253)

    # =========================================================================
    # DIAPOSITIVA 2: EL CONTEXTO
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    header = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_OSCURO
    header.line.fill.background()
    
    tx_header = slide2.shapes.add_textbox(Inches(0.75), Inches(0.15), Inches(11.83), Inches(0.7))
    p_h = tx_header.text_frame.paragraphs[0]
    p_h.text = "¿Por qué cambia la facturación en España?"
    p_h.font.bold = True
    p_h.font.size = Pt(24)
    p_h.font.color.rgb = BLANCO

    txBox = slide2.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "1. Fin de las facturas manuales o en Word/Excel sueltos"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = AZUL_OSCURO
    p.space_after = Pt(6)

    p_sub = tf.add_paragraph()
    p_sub.text = "   - Hacienda digitaliza los negocios para evitar la contabilidad paralela o en negro."
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = TEXTO_OSCURO
    p_sub.space_after = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "2. Mayor control fiscal para evitar el fraude de doble uso"
    p2.font.bold = True
    p2.font.size = Pt(20)
    p2.font.color.rgb = AZUL_OSCURO
    p2.space_after = Pt(6)

    p2_sub = tf.add_paragraph()
    p2_sub.text = "   - Todo el software de facturación debe estar homologado e impedir que se borren o modifiquen registros."
    p2_sub.font.size = Pt(15)
    p2_sub.font.color.rgb = TEXTO_OSCURO
    p2_sub.space_after = Pt(14)

    p3 = tf.add_paragraph()
    p3.text = "3. Reducción de la morosidad comercial (Plazos de Cobro)"
    p3.font.bold = True
    p3.font.size = Pt(20)
    p3.font.color.rgb = AZUL_OSCURO
    p3.space_after = Pt(6)

    p3_sub = tf.add_paragraph()
    p3_sub.text = "   - Se obligará a registrar cuándo se cobran las facturas reales para vigilar que se pague antes de 60 días."
    p3_sub.font.size = Pt(15)
    p3_sub.font.color.rgb = TEXTO_OSCURO

    # =========================================================================
    # DIAPOSITIVA 3: REGLAMENTO VERIFACTU (CON IMAGEN DEL ENCADENAMIENTO)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    header = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_OSCURO
    header.line.fill.background()
    
    tx_header = slide3.shapes.add_textbox(Inches(0.75), Inches(0.15), Inches(11.83), Inches(0.7))
    p_h = tx_header.text_frame.paragraphs[0]
    p_h.text = "Reglamento Veri*factu: La regla del bolígrafo de tinta"
    p_h.font.bold = True
    p_h.font.size = Pt(24)
    p_h.font.color.rgb = BLANCO

    # Columna izquierda: Conceptos
    txLeft = slide3.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(5.0))
    tfL = txLeft.text_frame
    tfL.word_wrap = True

    p = tfL.paragraphs[0]
    p.text = "Conceptos Clave de Veri*factu:"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = AZUL_OSCURO
    p.space_after = Pt(14)

    bullet1 = tfL.add_paragraph()
    bullet1.text = "• Inalterabilidad absoluta:"
    bullet1.font.bold = True
    bullet1.font.size = Pt(15)
    bullet1.font.color.rgb = TEXTO_OSCURO
    
    b1_sub = tfL.add_paragraph()
    b1_sub.text = "  No se pueden borrar facturas. Si hay un error, se genera una factura rectificativa de anulación."
    b1_sub.font.size = Pt(13)
    b1_sub.font.color.rgb = TEXTO_MUTED
    b1_sub.space_after = Pt(10)

    bullet2 = tfL.add_paragraph()
    bullet2.text = "• Encadenamiento criptográfico:"
    bullet2.font.bold = True
    bullet2.font.size = Pt(15)
    bullet2.font.color.rgb = TEXTO_OSCURO

    b2_sub = tfL.add_paragraph()
    b2_sub.text = "  Cada factura lleva cosida digitalmente la huella de la anterior. Alterar una factura antigua rompe toda la cadena."
    b2_sub.font.size = Pt(13)
    b2_sub.font.color.rgb = TEXTO_MUTED
    b2_sub.space_after = Pt(10)

    bullet3 = tfL.add_paragraph()
    bullet3.text = "• Código QR de Verificación:"
    bullet3.font.bold = True
    bullet3.font.size = Pt(15)
    bullet3.font.color.rgb = TEXTO_OSCURO

    b3_sub = tfL.add_paragraph()
    b3_sub.text = "  El cliente puede escanear la factura con su móvil y verificar en la sede web de la AEAT que es legal al instante."
    b3_sub.font.size = Pt(13)
    b3_sub.font.color.rgb = TEXTO_MUTED

    # Columna derecha: Imagen explicativa del encadenamiento
    boxRight = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.0))
    boxRight.fill.solid()
    boxRight.fill.fore_color.rgb = AZUL_CLARO
    boxRight.line.color.rgb = RGBColor(191, 219, 254)

    # Insertar Imagen del Encadenamiento
    try:
        slide3.shapes.add_picture(img_cadena, Inches(7.1), Inches(2.2), width=Inches(5.2), height=Inches(3.8))
        
        # Etiqueta de la imagen
        lbl = slide3.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.8), Inches(0.5))
        lbl.text_frame.paragraphs[0].text = "🔗 ESQUEMA VISUAL DEL ENCADENAMIENTO:"
        lbl.text_frame.paragraphs[0].font.bold = True
        lbl.text_frame.paragraphs[0].font.size = Pt(14)
        lbl.text_frame.paragraphs[0].font.color.rgb = AMBAR
        lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    except Exception as e:
        print("Error insertando imagen de cadena:", e)

    # =========================================================================
    # DIAPOSITIVA 4: TABLA COMPARATIVA (EL PORQUÉ DEL NUEVO SISTEMA)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    header = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_OSCURO
    header.line.fill.background()
    
    tx_header = slide4.shapes.add_textbox(Inches(0.75), Inches(0.15), Inches(11.83), Inches(0.7))
    p_h = tx_header.text_frame.paragraphs[0]
    p_h.text = "Comparativa: ¿Por qué hay que facturar así ahora?"
    p_h.font.bold = True
    p_h.font.size = Pt(24)
    p_h.font.color.rgb = BLANCO

    # Añadir Tabla Comparativa
    rows, cols = 5, 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.33)
    height = Inches(5.2)

    table_shape = slide4.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Anchos de columna
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(3.2)
    table.columns[2].width = Inches(3.7)
    table.columns[3].width = Inches(3.63)

    headers_txt = ["Aspecto", "Facturación Antigua (Manual / Word)", "Nueva Facturación (Ley Verifactu / Crea)", "¿Por qué se hace así ahora? (Motivo)"]
    data = [
        ["Borrar/Editar\nFacturas", "Se podía corregir un número o reescribir un concepto en un Word/Excel.", "Totalmente prohibido. Debe emitirse una factura rectificativa nueva.", "Para asegurar la transparencia y evitar que se oculten o eliminen ventas reales."],
        ["Formato del\nDocumento", "Papel físico impreso o PDF plano enviado por correo electrónico común.", "Ficheros electrónicos estructurados (XML / Facturae) o con firma y QR.", "Para permitir el cruce automático de datos fiscales y agilizar las auditorías."],
        ["Envío a la\nAdministración", "Se declaraba en resúmenes globales cada 3 meses (Modelo 303/390).", "Envío inmediato de cada registro de factura a la AEAT al facturar.", "Para evitar que los datos fiscales se manipulen tarde o se pierdan."],
        ["Control de\nCobros (B2B)", "Sin control. Retrasos y morosidad habituales sin consecuencias oficiales.", "Obligatoriedad de informar la fecha del cobro bancario real de cada factura.", "Para combatir la morosidad y proteger la liquidez de los autónomos."]
    ]

    # Dar formato a la cabecera
    for col_idx, text in enumerate(headers_txt):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_OSCURO
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = BLANCO
            paragraph.alignment = PP_ALIGN.CENTER

    # Rellenar datos
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_value
            cell.fill.solid()
            # Fila alterna
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = AZUL_CLARO
            else:
                cell.fill.fore_color.rgb = BLANCO
                
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11.5)
                paragraph.font.color.rgb = TEXTO_OSCURO
                # Columna de Aspecto y Por Qué en negrita suave
                if col_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = AZUL_OSCURO
                elif col_idx == 3:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = VERDE

    # =========================================================================
    # DIAPOSITIVA 5: LEY CREA Y CRECE
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    header = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_OSCURO
    header.line.fill.background()
    
    tx_header = slide5.shapes.add_textbox(Inches(0.75), Inches(0.15), Inches(11.83), Inches(0.7))
    p_h = tx_header.text_frame.paragraphs[0]
    p_h.text = "Ley Crea y Crece: El adiós definitivo al PDF e-mail"
    p_h.font.bold = True
    p_h.font.size = Pt(24)
    p_h.font.color.rgb = BLANCO

    txBox = slide5.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "La Factura Electrónica Obligatoria entre Negocios (B2B)"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = AZUL_OSCURO
    p.space_after = Pt(14)

    b1 = tf.add_paragraph()
    b1.text = "• Formato digital estructurado obligatorio:"
    b1.font.bold = True
    b1.font.size = Pt(16)
    b1.font.color.rgb = TEXTO_OSCURO

    b1_sub = tf.add_paragraph()
    b1_sub.text = "  Ya no vale mandar un PDF por email. Las facturas se deben emitir en ficheros XML (Facturae o UBL) firmados digitalmente, transmitidos directamente al sistema del cliente."
    b1_sub.font.size = Pt(14)
    b1_sub.font.color.rgb = TEXTO_MUTED
    b1_sub.space_after = Pt(12)

    b2 = tf.add_paragraph()
    b2.text = "• Obligatoriedad de informar los estados de pago:"
    b2.font.bold = True
    b2.font.size = Pt(16)
    b2.font.color.rgb = TEXTO_OSCURO

    b2_sub = tf.add_paragraph()
    b2_sub.text = "  Ambas partes deben reportar cuándo se recibe la factura y el día exacto en que se realiza el cobro bancario real."
    b2_sub.font.size = Pt(14)
    b2_sub.font.color.rgb = TEXTO_MUTED
    b2_sub.space_after = Pt(12)

    b3 = tf.add_paragraph()
    b3.text = "• Lucha contra la morosidad comercial:"
    b3.font.bold = True
    b3.font.size = Pt(16)
    b3.font.color.rgb = TEXTO_OSCURO

    b3_sub = tf.add_paragraph()
    b3_sub.text = "  Hacienda vigilará que las empresas paguen en el plazo legal (máximo 60 días). Quien incumpla no podrá acceder a subvenciones."
    b3_sub.font.size = Pt(14)
    b3_sub.font.color.rgb = TEXTO_MUTED

    # =========================================================================
    # DIAPOSITIVA 6: SOLUCIÓN FACTURAFÁCIL (CON IMAGEN DEL FLUJO DE TRABAJO)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    header = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_OSCURO
    header.line.fill.background()
    
    tx_header = slide6.shapes.add_textbox(Inches(0.75), Inches(0.15), Inches(11.83), Inches(0.7))
    p_h = tx_header.text_frame.paragraphs[0]
    p_h.text = "Nuestra Solución: FacturaFácil para Autónomos"
    p_h.font.bold = True
    p_h.font.size = Pt(24)
    p_h.font.color.rgb = BLANCO

    # Columna izquierda: Características
    txLeft = slide6.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(5.0))
    tfL = txLeft.text_frame
    tfL.word_wrap = True

    p = tfL.paragraphs[0]
    p.text = "¿Cómo te ayuda el programa?"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = AZUL_OSCURO
    p.space_after = Pt(14)

    bullet1 = tfL.add_paragraph()
    bullet1.text = "1. Fácil y sin tecnicismos:"
    bullet1.font.bold = True
    bullet1.font.size = Pt(15)
    bullet1.font.color.rgb = TEXTO_OSCURO
    
    b1_sub = tfL.add_paragraph()
    b1_sub.text = "   Solo introduces qué has vendido e importes. La numeración y cálculos se hacen solos."
    b1_sub.font.size = Pt(13)
    b1_sub.font.color.rgb = TEXTO_MUTED
    b1_sub.space_after = Pt(10)

    bullet2 = tfL.add_paragraph()
    bullet2.text = "2. Cumplimiento legal automático:"
    bullet2.font.bold = True
    bullet2.font.size = Pt(15)
    bullet2.font.color.rgb = TEXTO_OSCURO

    b2_sub = tfL.add_paragraph()
    b2_sub.text = "   Calcula la huella criptográfica de Veri*factu y genera el código QR sin que tengas que intervenir."
    b2_sub.font.size = Pt(13)
    b2_sub.font.color.rgb = TEXTO_MUTED
    b2_sub.space_after = Pt(10)

    bullet3 = tfL.add_paragraph()
    bullet3.text = "3. Copias en PDF de calidad:"
    bullet3.font.bold = True
    bullet3.font.size = Pt(15)
    bullet3.font.color.rgb = TEXTO_OSCURO

    b3_sub = tfL.add_paragraph()
    b3_sub.text = "   Con un clic puedes imprimir en papel o guardar en tu ordenador la factura en formato oficial A4."
    b3_sub.font.size = Pt(13)
    b3_sub.font.color.rgb = TEXTO_MUTED

    # Columna derecha: Imagen explicativa del flujo
    boxRight = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.0))
    boxRight.fill.solid()
    boxRight.fill.fore_color.rgb = AZUL_CLARO
    boxRight.line.color.rgb = RGBColor(191, 219, 254)

    # Insertar Imagen del Flujo
    try:
        slide6.shapes.add_picture(img_flujo, Inches(7.1), Inches(2.2), width=Inches(5.2), height=Inches(3.8))
        
        # Etiqueta de la imagen
        lbl = slide6.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.8), Inches(0.5))
        lbl.text_frame.paragraphs[0].text = "💼 PROCESO DE FACTURACIÓN EN 4 PASOS:"
        lbl.text_frame.paragraphs[0].font.bold = True
        lbl.text_frame.paragraphs[0].font.size = Pt(14)
        lbl.text_frame.paragraphs[0].font.color.rgb = AMBAR
        lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    except Exception as e:
        print("Error insertando imagen de flujo:", e)

    # =========================================================================
    # DIAPOSITIVA 7: HOJA DE RUTA
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    header = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_OSCURO
    header.line.fill.background()
    
    tx_header = slide7.shapes.add_textbox(Inches(0.75), Inches(0.15), Inches(11.83), Inches(0.7))
    p_h = tx_header.text_frame.paragraphs[0]
    p_h.text = "Hoja de ruta: Plan de acción sencillo"
    p_h.font.bold = True
    p_h.font.size = Pt(24)
    p_h.font.color.rgb = BLANCO

    txBox = slide7.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Cómo prepararse con tiempo y sin agobios"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = AZUL_OSCURO
    p.space_after = Pt(14)

    c1 = tf.add_paragraph()
    c1.text = "1. Practicar con la demo actual"
    c1.font.bold = True
    c1.font.size = Pt(16)
    c1.font.color.rgb = TEXTO_OSCURO

    c1_sub = tf.add_paragraph()
    c1_sub.text = "   Dedica unos minutos a crear clientes ficticios y hacer facturas simuladas en FacturaFácil para habituarte al cambio."
    c1_sub.font.size = Pt(15)
    c1_sub.font.color.rgb = TEXTO_MUTED
    c1_sub.space_after = Pt(10)

    c2 = tf.add_paragraph()
    c2.text = "2. Configurar los datos de tu empresa"
    c2.font.bold = True
    c2.font.size = Pt(16)
    c2.font.color.rgb = TEXTO_OSCURO

    c2_sub = tf.add_paragraph()
    c2_sub.text = "   Completa tu NIF y dirección en la pestaña de configuración para tener la herramienta lista para facturar en tu día a día."
    c2_sub.font.size = Pt(15)
    c2_sub.font.color.rgb = TEXTO_MUTED
    c2_sub.space_after = Pt(10)

    c3 = tf.add_paragraph()
    c3.text = "3. Estar listos antes de 2027"
    c3.font.bold = True
    c3.font.size = Pt(16)
    c3.font.color.rgb = TEXTO_OSCURO

    c3_sub = tf.add_paragraph()
    c3_sub.text = "   Disponemos de margen de tiempo de sobra. Adaptarse antes de las fechas límite nos garantiza una transición tranquila y segura."
    c3_sub.font.size = Pt(15)
    c3_sub.font.color.rgb = TEXTO_MUTED

    # Guardar presentación
    import os
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.abspath(os.path.join(script_dir, "..", "resources", "presentacion_hacienda_visual.pptx"))
    prs.save(output_path)
    print("Presentacion visual creada correctamente en: " + output_path)

if __name__ == "__main__":
    crear_presentacion()
